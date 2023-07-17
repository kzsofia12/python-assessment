from pptx import Presentation
from slide import Slide
import plotly.graph_objs as plot
import plotly.io as plot_io
import numpy as np
import csv
from pptx.util import Pt
import logging

# define constant size parameters
LEFT = Pt(120)
TOP = Pt(90)
WIDTH = Pt(500)
HEIGHT = Pt(400)

# setting log information
logging.basicConfig(filename='events.log',
                    filemode='a+',
                    format='%(asctime)s: %(levelname)s - %(message)s',
                    datefmt='%y-%m-%d %H:%M')
logger = logging.getLogger("log")
logger.setLevel(logging.INFO)


class SlideManager:
    # global variables
    output_filename = ""
    pres = Presentation()

    def __init__(self, output):
        self.output_filename = output

    # generate title slide
    def title_slide(self, properties: Slide):
        """Generate Title slide\n
        Contains a title and a sub-title
        """
        logger.info("title_slide: Start generate")
        # set title layout slide
        layout = self.pres.slide_layouts[0]

        # create slide
        new_slide = self.pres.slides.add_slide(layout)

        # add title and content
        try:
            new_slide.shapes.title.text = properties.title
            new_slide.placeholders[1].text = properties.content

        except Exception as e:
            logger.error("title_slide: text or content error: ", e)

        else:
            # save slide to the presentation
            self.pres.save(self.output_filename)
            logger.info("title_slide: Generate successfully")

    # add title and textbox slide
    def text_slide(self, properties: Slide):
        """Generate a slide with a long text.\n
        Contains a title and some text"""
        logger.info("text_slide: Start generate")

        # set title only slide
        layout = self.pres.slide_layouts[5]

        # create slide
        new_slide = self.pres.slides.add_slide(layout)

        # add title text
        try:
            new_slide.shapes.title.text = properties.title
        except Exception as e:
            logger.error("text_slide: error in title: ", e)

        # create a new textbox
        textbox = new_slide.shapes.add_textbox(LEFT, TOP, WIDTH, HEIGHT)

        # add a new paragraph to textbox
        text_frame = textbox.text_frame
        paragraph = text_frame.add_paragraph()

        # add text to textbox
        try:
            paragraph.text = properties.content
        except Exception as e:
            logger.error("text_slide: error in content: ", e)

        # save slide to presentation
        self.pres.save(self.output_filename)
        logger.info("text_slide: Generate successfully")

    def list_slide(self, properties: Slide):
        """Generate a slide with a list\n
        Contains a title and a bullet list"""
        logger.info("list_slide: Start generate")

        # set title and content slide
        layout = self.pres.slide_layouts[1]

        # create slide
        new_slide = self.pres.slides.add_slide(layout)

        # add title text
        try:
            new_slide.shapes.title.text = properties.title
        except Exception as e:
            logger.error("list_slide: error in title: ", e)

        # add a new placeholder
        place = new_slide.placeholders[1]
        tf = place.text_frame

        # go through the list elements
        try:
            for text in properties.content:
                par = tf.add_paragraph()

                # set the list level
                par.level = text['level']

                # set list text
                par.text = text['text']

        except Exception as e:
            logger.error("list_slide: error in content: ", e)

        # save slide to presentation
        self.pres.save(self.output_filename)
        logger.info("list_slide: Generate successfully")

    def image_slide(self, properties: Slide):
        """Generate a slide with an image\n
            Contains a title and an image"""
        logger.info("image_slide: Start generate")

        # set title only slide
        layout = self.pres.slide_layouts[5]

        # create slide
        new_slide = self.pres.slides.add_slide(layout)

        try:
            # add title text to slide
            new_slide.shapes.title.text = properties.title
        except Exception as e:
            logger.error("image_slide: error in title: ", e)

        try:
            # add picture to slide
            new_slide.shapes.add_picture(properties.content, LEFT, TOP, WIDTH, HEIGHT)
        except FileNotFoundError:
            logger.error("image_slide: Image not found")
        except Exception as e:
            logger.error("image_slide: error in content: ", e)

        # save slide to presentation
        self.pres.save(self.output_filename)
        logger.info("image_slide: Generate successfully")

    def plot_slide(self, properties: Slide):
        """Generate a slide with a plot\n
        Create a line chart from coordinates and insert into the slide"""
        logger.info("plot_slide: Start generate")

        try:
            # read coordinates from file
            with open(properties.content, newline='') as file:
                reader = csv.reader(file, delimiter=';')
                points = list(reader)

        except FileNotFoundError:
            logger.error("plot_slide: ",properties.content," not found")
        except Exception as e:
            logger.error("plot_slide: error in content: ", e)
        else:
            # modify 4 2x1 array to 2 4x1 array
            # First array contain the x coordinates, second contain the y coordinates
            coordinates = np.transpose(points)
            try:
                # set the axis titles
                chart_layout = plot.Layout(xaxis_title=properties.configuration['x-label'],
                                           yaxis_title=properties.configuration['y-label'])

                # create linechart
                chart = plot.Figure(data=plot.Scatter(x=coordinates[0], y=coordinates[1]), layout=chart_layout)

                # import the chart to png
                plot_io.write_image(chart, 'chart.png')
            except Exception as e:
                logger.error("plot_slide: error in create chart: ", e)
            else:
                # define only title slide
                layout = self.pres.slide_layouts[5]

                # create slide
                new_slide = self.pres.slides.add_slide(layout)
                try:
                    # add title
                    new_slide.shapes.title.text = properties.title
                except Exception as e:
                    logger.error("plot_slide: error in title: ", e)

                try:
                    # add chart picture
                    new_slide.shapes.add_picture('chart.png', LEFT, TOP, WIDTH, HEIGHT)
                except Exception as e:
                    logger.error("plot_slide: error when add the chart image: ", e)

                # save slide to presentation
                self.pres.save(self.output_filename)
                logger.info("plot_slide: Generate successfully")
